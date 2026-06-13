#!/usr/bin/env python3
"""
Comprehensive script to update the Phase8 mini project PPTX:
1. Update slides 2, 3, 5, 6, 7, 8 with actual project content
2. Add new slides after slide 68 for subtopic 4 content
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree

# ─── Paths ───
BASE = "/Volumes/Corsair EX300U Media/00_work_out/02_ing/Phase8_mini"
PPTX_PATH = os.path.join(BASE, "ppt", "Phase8_ 미니 프로젝트.pptx")
IMG_DIR = os.path.join(BASE, "outputs", "figures", "subtopic4_reorder_clustering")

prs = Presentation(PPTX_PATH)

# ─── Constants ───
SLIDE_WIDTH = prs.slide_width   # 24384000 EMU
SLIDE_HEIGHT = prs.slide_height # 13716000 EMU

COLOR_DARK = RGBColor(0x2F, 0x2F, 0x2F)
COLOR_GRAY = RGBColor(0x78, 0x78, 0x78)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x11, 0x11, 0x11)
COLOR_LIGHT_GRAY = RGBColor(0xB7, 0xB7, 0xB7)

FONT_NAME = "Arial"


# ═══════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════

def set_text_props(run, font_name=FONT_NAME, font_size=None, color=None, bold=None):
    """Set run font properties."""
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = font_size
    if color:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold


def clear_and_set_text(text_frame, text, font_name=FONT_NAME, font_size=None, color=None, bold=None):
    """Clear all paragraphs in text_frame and set single paragraph text."""
    # Clear existing paragraphs
    for i in range(len(text_frame.paragraphs) - 1, 0, -1):
        p = text_frame.paragraphs[i]._p
        p.getparent().remove(p)
    # Set first paragraph
    para = text_frame.paragraphs[0]
    para.clear()
    run = para.add_run()
    run.text = text
    set_text_props(run, font_name, font_size, color, bold)
    return para


def set_multiline_text(text_frame, lines, font_name=FONT_NAME, font_size=None, color=None, bold=None, alignment=None):
    """Set multiple lines in text_frame, each as a separate paragraph."""
    # Clear existing paragraphs
    for i in range(len(text_frame.paragraphs) - 1, 0, -1):
        p = text_frame.paragraphs[i]._p
        p.getparent().remove(p)

    for idx, line in enumerate(lines):
        if idx == 0:
            para = text_frame.paragraphs[0]
            para.clear()
        else:
            para = text_frame.add_paragraph()
        run = para.add_run()
        run.text = line
        set_text_props(run, font_name, font_size, color, bold)
        if alignment is not None:
            para.alignment = alignment


def find_shape_by_name(slide, name):
    """Find shape by name in slide."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def add_textbox(slide, left, top, width, height, text, font_name=FONT_NAME, font_size=Pt(30), color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    set_text_props(run, font_name, font_size, color, bold)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_name=FONT_NAME, font_size=Pt(24), color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=None):
    """Add a multi-line text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = alignment
        if line_spacing:
            para.space_after = Pt(line_spacing)
        run = para.add_run()
        run.text = line
        set_text_props(run, font_name, font_size, color, bold)
    return txBox


def add_image_centered(slide, img_path, top, max_width, max_height):
    """Add image centered horizontally, fitting within max dimensions."""
    if not os.path.exists(img_path):
        print(f"WARNING: Image not found: {img_path}")
        return None
    from PIL import Image
    img = Image.open(img_path)
    img_w, img_h = img.size
    # Calculate scale
    scale_w = max_width / img_w
    scale_h = max_height / img_h
    scale = min(scale_w, scale_h)
    final_w = int(img_w * scale)
    final_h = int(img_h * scale)
    # Center horizontally
    left = (SLIDE_WIDTH - Emu(final_w * 914400 // 96)) // 2  # rough conversion
    # Actually use Inches-based calculation
    # EMU per pixel at 96 DPI = 914400/96 = 9525
    emu_w = int(final_w * 9525)
    emu_h = int(final_h * 9525)
    left = (SLIDE_WIDTH - emu_w) // 2
    pic = slide.shapes.add_picture(img_path, left, top, emu_w, emu_h)
    return pic


def add_new_blank_slide(prs, after_index):
    """Add a new blank slide after the given index."""
    blank_layout = prs.slide_layouts[0]  # BLANK layout
    slide = prs.slides.add_slide(blank_layout)
    # Move slide to correct position
    # python-pptx adds slides at the end, we need to reorder
    slides_element = prs.slides._sldIdLst
    slide_id_list = list(slides_element)
    # The new slide is the last one
    new_slide_element = slide_id_list[-1]
    # Remove from end
    slides_element.remove(new_slide_element)
    # Insert after the target index
    if after_index + 1 < len(slide_id_list) - 1:
        ref_element = slide_id_list[after_index + 1]
        slides_element.insert(slides_element.index(ref_element), new_slide_element)
    else:
        slides_element.append(new_slide_element)
    return slide


def add_slide_header(slide, section_number, section_name):
    """Add standard header to slide (logo area text + section number)."""
    # Company name (top-left) - skip logo image, just add text
    add_textbox(slide, Emu(2146300), Emu(1231900), Emu(2108100), Emu(533400),
                "미리컴퍼니", FONT_NAME, Pt(27), COLOR_DARK, False)
    # Section label (top-right)
    add_textbox(slide, Emu(20091400), Emu(1282700), Emu(3009900), Emu(406500),
                section_name, FONT_NAME, Pt(20), COLOR_DARK, False, PP_ALIGN.RIGHT)
    # Footer
    add_textbox(slide, Emu(13093700), Emu(12407900), Emu(10007700), Emu(355500),
                "www.miricompany.com", FONT_NAME, Pt(18), COLOR_DARK, False)


def add_slide_title(slide, title_text, font_size=Pt(48)):
    """Add main title to a new slide."""
    add_textbox(slide, Emu(1905000), Emu(2044700), Emu(20000000), Emu(1359000),
                title_text, FONT_NAME, font_size, COLOR_DARK, True, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════
# PART 1: UPDATE EXISTING SLIDES (2, 3, 5, 6, 7, 8)
# ═══════════════════════════════════════════════════════════

print("=" * 60)
print("PART 1: Updating existing slides")
print("=" * 60)

# ─── SLIDE 2: 목차 (Table of Contents) ───
print("Updating Slide 2 (목차)...")
slide2 = prs.slides[1]

# Update section header label (top-right)
shape = find_shape_by_name(slide2, "Google Shape;187;g3cfba76c0e1_0_0")
if shape and shape.has_text_frame:
    clear_and_set_text(shape.text_frame, "00. 목차", FONT_NAME, Pt(20), COLOR_DARK)

# Update TOC items - 01 through 05
# 01: "기획 배경 및 목표" -> "팀 빌딩 & 기획 배경"
shape188 = find_shape_by_name(slide2, "Google Shape;188;g3cfba76c0e1_0_0")
if shape188 and shape188.has_text_frame:
    clear_and_set_text(shape188.text_frame, "기획 배경 및 목적", FONT_NAME, Pt(34), COLOR_DARK)

# 01 sub-items
shape203 = find_shape_by_name(slide2, "Google Shape;203;g3cfba76c0e1_0_0")
if shape203 and shape203.has_text_frame:
    clear_and_set_text(shape203.text_frame, "- 배경 및 문제 인식", FONT_NAME, Pt(22), COLOR_DARK)
shape204 = find_shape_by_name(slide2, "Google Shape;204;g3cfba76c0e1_0_0")
if shape204 and shape204.has_text_frame:
    clear_and_set_text(shape204.text_frame, "- 기획 의도 및 목적", FONT_NAME, Pt(22), COLOR_DARK)

# 02: "시장 및 고객 분석" -> "프로젝트 팀 구성 및 역할"
shape191 = find_shape_by_name(slide2, "Google Shape;191;g3cfba76c0e1_0_0")
if shape191 and shape191.has_text_frame:
    clear_and_set_text(shape191.text_frame, "프로젝트 팀 구성 및 역할", FONT_NAME, Pt(34), COLOR_DARK)
shape207 = find_shape_by_name(slide2, "Google Shape;207;g3cfba76c0e1_0_0")
if shape207 and shape207.has_text_frame:
    clear_and_set_text(shape207.text_frame, "- 굿핏(good fit) 4인 팀", FONT_NAME, Pt(22), COLOR_DARK)
shape208 = find_shape_by_name(slide2, "Google Shape;208;g3cfba76c0e1_0_0")
if shape208 and shape208.has_text_frame:
    clear_and_set_text(shape208.text_frame, "- 역할 분담 및 소주제 배정", FONT_NAME, Pt(22), COLOR_DARK)

# 03: "마케팅 전략 및 컨셉" -> "프로젝트 수행 절차 및 방법"
shape194 = find_shape_by_name(slide2, "Google Shape;194;g3cfba76c0e1_0_0")
if shape194 and shape194.has_text_frame:
    clear_and_set_text(shape194.text_frame, "프로젝트 수행 절차 및 방법", FONT_NAME, Pt(34), COLOR_DARK)
shape211 = find_shape_by_name(slide2, "Google Shape;211;g3cfba76c0e1_0_0")
if shape211 and shape211.has_text_frame:
    clear_and_set_text(shape211.text_frame, "- 기(起)-승(承)-전(轉)-결(結) 흐름", FONT_NAME, Pt(22), COLOR_DARK)
shape212 = find_shape_by_name(slide2, "Google Shape;212;g3cfba76c0e1_0_0")
if shape212 and shape212.has_text_frame:
    clear_and_set_text(shape212.text_frame, "- 4개 소주제 구성", FONT_NAME, Pt(22), COLOR_DARK)

# 04: "실행 계획" -> "사용 데이터 & 기대효과"
shape197 = find_shape_by_name(slide2, "Google Shape;197;g3cfba76c0e1_0_0")
if shape197 and shape197.has_text_frame:
    clear_and_set_text(shape197.text_frame, "사용 데이터 & 기대효과", FONT_NAME, Pt(34), COLOR_DARK)
shape205 = find_shape_by_name(slide2, "Google Shape;205;g3cfba76c0e1_0_0")
if shape205 and shape205.has_text_frame:
    clear_and_set_text(shape205.text_frame, "- Kaggle E-Grocery 1,000x37", FONT_NAME, Pt(22), COLOR_DARK)
shape206 = find_shape_by_name(slide2, "Google Shape;206;g3cfba76c0e1_0_0")
if shape206 and shape206.has_text_frame:
    clear_and_set_text(shape206.text_frame, "- 5대 인사이트 도출 목표", FONT_NAME, Pt(22), COLOR_DARK)

# 05: "예산 및 성과 관리" -> "소주제 1~4 분석 결과"
shape200 = find_shape_by_name(slide2, "Google Shape;200;g3cfba76c0e1_0_0")
if shape200 and shape200.has_text_frame:
    clear_and_set_text(shape200.text_frame, "소주제 1~4 분석 결과", FONT_NAME, Pt(34), COLOR_DARK)
shape209 = find_shape_by_name(slide2, "Google Shape;209;g3cfba76c0e1_0_0")
if shape209 and shape209.has_text_frame:
    clear_and_set_text(shape209.text_frame, "- 분류/회귀/이진분류/클러스터링", FONT_NAME, Pt(22), COLOR_DARK)
shape210 = find_shape_by_name(slide2, "Google Shape;210;g3cfba76c0e1_0_0")
if shape210 and shape210.has_text_frame:
    clear_and_set_text(shape210.text_frame, "- 통합 재고 관리 전략 도출", FONT_NAME, Pt(22), COLOR_DARK)


# ─── SLIDE 3: 기획 배경 및 목적 ───
print("Updating Slide 3 (기획 배경 및 목적)...")
slide3 = prs.slides[2]

# Section label
shape224 = find_shape_by_name(slide3, "Google Shape;224;g3cfba76c0e1_0_43")
if shape224 and shape224.has_text_frame:
    clear_and_set_text(shape224.text_frame, "01. 기획 배경 및 목적", FONT_NAME, Pt(20), COLOR_DARK)

# Subtitle / intro text
shape227 = find_shape_by_name(slide3, "Google Shape;227;g3cfba76c0e1_0_43")
if shape227 and shape227.has_text_frame:
    clear_and_set_text(shape227.text_frame,
        "데이터 기반 재고 관리 최적화로 식료품 유통의 전 주기를 진단합니다.",
        FONT_NAME, Pt(22), COLOR_DARK)

# Background section - "기획 배경" content
shape234 = find_shape_by_name(slide3, "Google Shape;234;g3cfba76c0e1_0_43")
if shape234 and shape234.has_text_frame:
    set_multiline_text(shape234.text_frame, [
        "식료품 유통업에서 재고 관리의 비효율은 곧 비용 손실로 직결",
        "과잉 재고 → 유통기한 초과 폐기 / 재고 부족 → 판매 기회 상실",
        "국내 중소기업 ERP 도입률 16.3% (대기업 95.5%)"
    ], FONT_NAME, Pt(27), COLOR_DARK)

# Purpose section - "기획 목적" content
shape239 = find_shape_by_name(slide3, "Google Shape;239;g3cfba76c0e1_0_43")
if shape239 and shape239.has_text_frame:
    set_multiline_text(shape239.text_frame, [
        "분류/회귀/비지도학습을 아우르는 4개 소주제로 재고 관리 전 주기 분석",
        "머신러닝 기반 재고 관리 의사결정 지원 시스템 구현",
        "상태 진단 → 수요 예측 → 위험 감지 → 전략 수립 사이클 완성"
    ], FONT_NAME, Pt(27), COLOR_DARK)


# ─── SLIDE 5: 기대효과 ───
print("Updating Slide 5 (기대효과)...")
slide5 = prs.slides[4]

# Section label
shape270 = find_shape_by_name(slide5, "Google Shape;270;g3cfba76c0e1_0_121")
if shape270 and shape270.has_text_frame:
    clear_and_set_text(shape270.text_frame, "04. 기대효과", FONT_NAME, Pt(20), COLOR_DARK)

# Effect 1 title: "브랜드 인지도 향상" -> "상태 진단 자동화"
shape287 = find_shape_by_name(slide5, "Google Shape;287;g3cfba76c0e1_0_121")
if shape287 and shape287.has_text_frame:
    clear_and_set_text(shape287.text_frame, "상태 진단 자동화", FONT_NAME, Pt(31), COLOR_WHITE)
# Effect 1 description
shape288 = find_shape_by_name(slide5, "Google Shape;288;g3cfba76c0e1_0_121")
if shape288 and shape288.has_text_frame:
    clear_and_set_text(shape288.text_frame,
        "재고 상태(In Stock/Low Stock/Expiring Soon)를 99% 정확도로 자동 분류하여 실시간 모니터링 가능",
        FONT_NAME, Pt(24), COLOR_DARK)

# Effect 2 title: "매출 증대 및 전환율 향상" -> "수요 예측 정확도"
shape281 = find_shape_by_name(slide5, "Google Shape;281;g3cfba76c0e1_0_121")
if shape281 and shape281.has_text_frame:
    clear_and_set_text(shape281.text_frame, "수요 예측 정확도 향상", FONT_NAME, Pt(31), COLOR_WHITE)
shape282 = find_shape_by_name(slide5, "Google Shape;282;g3cfba76c0e1_0_121")
if shape282 and shape282.has_text_frame:
    clear_and_set_text(shape282.text_frame,
        "일일 판매량 예측 R²=0.95로 발주 타이밍 및 수량 최적화 지원, 과재고/품절 동시 예방",
        FONT_NAME, Pt(24), COLOR_DARK)
shape283 = find_shape_by_name(slide5, "Google Shape;283;g3cfba76c0e1_0_121")
if shape283 and shape283.has_text_frame:
    clear_and_set_text(shape283.text_frame,
        "폐기 위험 제품 사전 식별(F1=0.99)로 False Negative 0건 달성, 폐기 비용 절감 기대",
        FONT_NAME, Pt(24), COLOR_DARK)

# Effect 3 title: "고객 참여 데이터 확보" -> "군집별 발주 전략"
shape292 = find_shape_by_name(slide5, "Google Shape;292;g3cfba76c0e1_0_121")
if shape292 and shape292.has_text_frame:
    clear_and_set_text(shape292.text_frame, "군집별 발주 전략 수립", FONT_NAME, Pt(31), COLOR_WHITE)
shape293 = find_shape_by_name(slide5, "Google Shape;293;g3cfba76c0e1_0_121")
if shape293 and shape293.has_text_frame:
    clear_and_set_text(shape293.text_frame,
        "K-Means 클러스터링으로 제품을 2개 군집으로 분류, EOQ 기반 차별화된 발주 전략 제시",
        FONT_NAME, Pt(24), COLOR_DARK)


# ─── SLIDE 6: 프로젝트 팀 구성 및 역할 ───
print("Updating Slide 6 (프로젝트 팀 구성 및 역할)...")
slide6 = prs.slides[5]

# Section label
shape302 = find_shape_by_name(slide6, "Google Shape;302;g3cfba76c0e1_0_107")
if shape302 and shape302.has_text_frame:
    clear_and_set_text(shape302.text_frame, "00. 팀 빌딩", FONT_NAME, Pt(20), COLOR_DARK)

# Update table (Shape 308) with team roles
shape308 = find_shape_by_name(slide6, "Google Shape;308;g3cfba76c0e1_0_107")
if shape308 and shape308.has_table:
    table = shape308.table
    # Header row
    cells_data = [
        ["이름", "역할", "담당 소주제", "주요 업무"],
        ["박준영", "팀장", "소주제 4", "최적 발주 전략 클러스터링\n소주제 통합 / EOQ / PM"],
        ["이현아", "팀원", "소주제 2", "일일 판매량 예측 (회귀)\nSHAP 분석"],
        ["권효중", "팀원", "소주제 3", "폐기 위험도 예측 (이진 분류)\n2x2 실험 설계"],
        ["정이랑", "팀원", "소주제 1", "재고 상태 분류 (다중 분류)\n기획 의도 제안"],
    ]
    for r_idx, row_data in enumerate(cells_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            # Clear cell and set text
            for para in cell.text_frame.paragraphs:
                para.clear()
            cell.text_frame.paragraphs[0].text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = cell_text
            run.font.name = FONT_NAME
            run.font.size = Pt(16)
            run.font.color.rgb = COLOR_DARK
            if r_idx == 0:
                run.font.bold = True
                run.font.color.rgb = COLOR_WHITE


# ─── SLIDE 7: 프로젝트 수행 절차 및 방법 ───
print("Updating Slide 7 (프로젝트 수행 절차 및 방법)...")
slide7 = prs.slides[6]

# Section label
shape316 = find_shape_by_name(slide7, "Google Shape;316;g3cfba76c0e1_0_177")
if shape316 and shape316.has_text_frame:
    clear_and_set_text(shape316.text_frame, "00. 수행 절차", FONT_NAME, Pt(20), COLOR_DARK)

# Update phase names (기획 -> 기(起), 준비 -> 승(承), 실행 -> 전(轉), 사후관리 -> 결(結), 최적화 -> 통합)
shape331 = find_shape_by_name(slide7, "Google Shape;331;g3cfba76c0e1_0_177")
if shape331 and shape331.has_text_frame:
    clear_and_set_text(shape331.text_frame, "기(起) 시작", FONT_NAME, Pt(26), COLOR_DARK)
shape336 = find_shape_by_name(slide7, "Google Shape;336;g3cfba76c0e1_0_177")
if shape336 and shape336.has_text_frame:
    clear_and_set_text(shape336.text_frame, "승(承) 전개", FONT_NAME, Pt(26), COLOR_DARK)
shape337 = find_shape_by_name(slide7, "Google Shape;337;g3cfba76c0e1_0_177")
if shape337 and shape337.has_text_frame:
    clear_and_set_text(shape337.text_frame, "전(轉) 발견", FONT_NAME, Pt(26), COLOR_DARK)
shape338 = find_shape_by_name(slide7, "Google Shape;338;g3cfba76c0e1_0_177")
if shape338 and shape338.has_text_frame:
    clear_and_set_text(shape338.text_frame, "결(結) 평가", FONT_NAME, Pt(26), COLOR_DARK)
shape339 = find_shape_by_name(slide7, "Google Shape;339;g3cfba76c0e1_0_177")
if shape339 and shape339.has_text_frame:
    clear_and_set_text(shape339.text_frame, "통합 전략", FONT_NAME, Pt(26), COLOR_DARK)

# Update date range
shape362 = find_shape_by_name(slide7, "Google Shape;362;g3cfba76c0e1_0_177")
if shape362 and shape362.has_text_frame:
    set_multiline_text(shape362.text_frame, [
        "2026",
        "03. 02 - 03. 16"
    ], FONT_NAME, Pt(17), RGBColor(0xC2, 0xC2, 0xC2))

# Phase dates
date_shapes = {
    "Google Shape;353;g3cfba76c0e1_0_177": "03. 02 - 03. 05",
    "Google Shape;354;g3cfba76c0e1_0_177": "03. 05 - 03. 09",
    "Google Shape;355;g3cfba76c0e1_0_177": "03. 09 - 03. 13",
    "Google Shape;356;g3cfba76c0e1_0_177": "03. 13 - 03. 15",
    "Google Shape;357;g3cfba76c0e1_0_177": "03. 15 - 03. 16",
}
for sname, date_text in date_shapes.items():
    s = find_shape_by_name(slide7, sname)
    if s and s.has_text_frame:
        clear_and_set_text(s.text_frame, date_text, FONT_NAME, Pt(12), COLOR_DARK)

# Update row labels
row_labels = {
    "Google Shape;332;g3cfba76c0e1_0_177": "데이터 전처리",
    "Google Shape;333;g3cfba76c0e1_0_177": "EDA & 누수 진단",
    "Google Shape;334;g3cfba76c0e1_0_177": "모델 학습 & 튜닝",
    "Google Shape;335;g3cfba76c0e1_0_177": "결과 분석 & 통합",
}
for sname, label in row_labels.items():
    s = find_shape_by_name(slide7, sname)
    if s and s.has_text_frame:
        clear_and_set_text(s.text_frame, label, FONT_NAME, Pt(22), COLOR_DARK)

# Row dates
row_dates = {
    "Google Shape;358;g3cfba76c0e1_0_177": "03. 02 - 03. 05",
    "Google Shape;359;g3cfba76c0e1_0_177": "03. 05 - 03. 09",
    "Google Shape;360;g3cfba76c0e1_0_177": "03. 09 - 03. 13",
    "Google Shape;361;g3cfba76c0e1_0_177": "03. 13 - 03. 16",
}
for sname, date_text in row_dates.items():
    s = find_shape_by_name(slide7, sname)
    if s and s.has_text_frame:
        clear_and_set_text(s.text_frame, date_text, FONT_NAME, Pt(12), COLOR_DARK)

# Update Gantt bar labels
bar_labels = {
    "Google Shape;341;g3cfba76c0e1_0_177": "주제 선정 / 데이터 수집 / 전처리",
    "Google Shape;344;g3cfba76c0e1_0_177": "EDA / 누수 진단 / 피처 엔지니어링",
    "Google Shape;346;g3cfba76c0e1_0_177": "모델 학습 / 하이퍼파라미터 튜닝 / 검증",
    "Google Shape;349;g3cfba76c0e1_0_177": "Feature Importance / SHAP / 결과 해석",
    "Google Shape;351;g3cfba76c0e1_0_177": "소주제 통합 / 발표 자료 / 회고",
}
for sname, label in bar_labels.items():
    s = find_shape_by_name(slide7, sname)
    if s and s.has_text_frame:
        clear_and_set_text(s.text_frame, label, FONT_NAME, Pt(19), COLOR_WHITE)

# Feedback labels
fb_labels = {
    "Google Shape;342;g3cfba76c0e1_0_177": "소주제 구성 확정 및 피드백",
    "Google Shape;347;g3cfba76c0e1_0_177": "분석 결과 공유 및 피드백",
    "Google Shape;352;g3cfba76c0e1_0_177": "모델 성능 검증 및 피드백",
}
for sname, label in fb_labels.items():
    s = find_shape_by_name(slide7, sname)
    if s and s.has_text_frame:
        clear_and_set_text(s.text_frame, label, FONT_NAME, Pt(12), COLOR_DARK)


# ─── SLIDE 8: 사용 데이터 ───
print("Updating Slide 8 (사용 데이터)...")
slide8 = prs.slides[7]

# Section label
shape371 = find_shape_by_name(slide8, "Google Shape;371;g3cfba76c0e1_0_153")
if shape371 and shape371.has_text_frame:
    clear_and_set_text(shape371.text_frame, "01. 사용 데이터", FONT_NAME, Pt(20), COLOR_DARK)

# URL text
shape382 = find_shape_by_name(slide8, "Google Shape;382;g3cfba76c0e1_0_153")
if shape382 and shape382.has_text_frame:
    clear_and_set_text(shape382.text_frame,
        "https://www.kaggle.com/datasets/mustofaahmad/inventory-management-grocery-industry",
        None, Pt(35), COLOR_LIGHT_GRAY)

# Data description title
shape384 = find_shape_by_name(slide8, "Google Shape;384;g3cfba76c0e1_0_153")
if shape384 and shape384.has_text_frame:
    clear_and_set_text(shape384.text_frame,
        "Kaggle - Inventory Management E-Grocery",
        FONT_NAME, Pt(30), COLOR_DARK, True)

# Data description body
shape383 = find_shape_by_name(slide8, "Google Shape;383;g3cfba76c0e1_0_153")
if shape383 and shape383.has_text_frame:
    set_multiline_text(shape383.text_frame, [
        "규모: 1,000행 x 37열 (SKU 단위)",
        "출처: 인도네시아 E-Grocery 운영 데이터",
        "카테고리: 10종 (Pantry, Beverages, Dairy 등)",
        "ABC 등급: A(200건), B(300건), C(500건)",
        "상태 레이블: In Stock / Low Stock / Out of Stock / Expiring Soon",
        "특이사항: 인도네시아 로케일(쉼표 소수점, $ 기호) 전처리 필요",
    ], FONT_NAME, Pt(27), COLOR_DARK)


# ═══════════════════════════════════════════════════════════
# PART 2: ADD NEW SLIDES AFTER SLIDE 68 (subtopic4)
# ═══════════════════════════════════════════════════════════

print("\n" + "=" * 60)
print("PART 2: Adding new slides for subtopic 4")
print("=" * 60)

# We'll add slides at the end and then reorder.
# Actually, since there are 68 slides and we want to add after 68 (the last slide),
# we can simply append new slides.

current_total = len(prs.slides)
print(f"Current total slides: {current_total}")

SECTION_LABEL = "소주제 4. 최적 발주 전략"


def create_new_slide():
    """Create a new blank slide at the end."""
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    return slide


def setup_slide(slide, header_text=SECTION_LABEL):
    """Add standard elements to a new slide."""
    add_slide_header(slide, "", header_text)


# ─── SLIDE S4-1: 소주제 4 표지 ───
print("  Adding S4-1: 소주제 4 표지...")
s4_1 = create_new_slide()
setup_slide(s4_1)
# Big title
add_textbox(s4_1, Emu(1905000), Emu(3500000), Emu(20000000), Emu(2500000),
            "소주제 4", FONT_NAME, Pt(68), COLOR_DARK, True)
add_textbox(s4_1, Emu(1905000), Emu(5800000), Emu(20000000), Emu(2000000),
            "최적 발주 전략 클러스터링", FONT_NAME, Pt(48), COLOR_DARK, False)
add_multiline_textbox(s4_1, Emu(1905000), Emu(8000000), Emu(20000000), Emu(2500000), [
    "Regression + K-Means Clustering + EOQ Simulation",
    "담당: 박준영 (팀장)",
], FONT_NAME, Pt(28), COLOR_GRAY)


# ─── SLIDE S4-2: 분석 배경 & 3단계 구조 ───
print("  Adding S4-2: 분석 배경 & 3단계 구조...")
s4_2 = create_new_slide()
setup_slide(s4_2)
add_slide_title(s4_2, "분석 배경 & 3단계 구조 (Phase A → B → C)")

add_multiline_textbox(s4_2, Emu(1905000), Emu(3600000), Emu(20000000), Emu(1500000), [
    "1,000개 이상의 제품을 동일 방식으로 발주하는 것은 비효율적",
    "제품 특성에 따라 발주 전략을 차별화해야 함",
    "소주제 1~3이 지도학습이라면, 소주제 4는 비지도학습으로 군집화",
], FONT_NAME, Pt(26), COLOR_DARK)

# 3-Phase table as text
add_multiline_textbox(s4_2, Emu(1905000), Emu(5800000), Emu(20000000), Emu(6000000), [
    "Phase A  |  재고 보유일수(DOI) 예측  |  회귀 (LR, RF, XGBoost)",
    "              →  DOI에 영향을 미치는 핵심 피처 도출",
    "",
    "Phase B  |  제품 군집화  |  K-Means Clustering",
    "              →  유사 제품 그룹 발견, 차별화 전략의 근거",
    "",
    "Phase C  |  발주량 최적화  |  EOQ 시뮬레이션",
    "              →  군집별 구체적 발주량 및 비용 산출",
    "",
    "세 단계가 연결되어야 '분석 → 이해 → 행동'의 완전한 흐름이 완성됨",
], FONT_NAME, Pt(24), COLOR_DARK)


# ─── SLIDE S4-3: EDA (DOI 분포 & 상관관계) ───
print("  Adding S4-3: EDA...")
s4_3 = create_new_slide()
setup_slide(s4_3)
add_slide_title(s4_3, "EDA: DOI 분포 및 피처 상관관계")

# DOI distribution image
img_path = os.path.join(IMG_DIR, "doi_distribution_s4.png")
if os.path.exists(img_path):
    pic = s4_3.shapes.add_picture(img_path, Emu(1200000), Emu(3600000), Emu(10500000), Emu(7500000))

# Correlation heatmap
img_path2 = os.path.join(IMG_DIR, "correlation_heatmap_s4.png")
if os.path.exists(img_path2):
    pic2 = s4_3.shapes.add_picture(img_path2, Emu(12500000), Emu(3600000), Emu(10500000), Emu(7500000))

add_multiline_textbox(s4_3, Emu(1200000), Emu(11400000), Emu(22000000), Emu(1000000), [
    "DOI 평균 ~10일(0~22일 분포)  |  QOH↔DOI +0.60  |  ADS↔DOI -0.38  |  카테고리 간 DOI 차이 미미",
], FONT_NAME, Pt(20), COLOR_GRAY)


# ─── SLIDE S4-4: EDA (카테고리별 DOI) ───
print("  Adding S4-4: 카테고리별 DOI...")
s4_4 = create_new_slide()
setup_slide(s4_4)
add_slide_title(s4_4, "EDA: 카테고리별 DOI 분포")

img_path = os.path.join(IMG_DIR, "category_doi_boxplot_s4.png")
if os.path.exists(img_path):
    pic = s4_4.shapes.add_picture(img_path, Emu(3000000), Emu(3400000), Emu(18000000), Emu(8500000))

add_multiline_textbox(s4_4, Emu(1905000), Emu(12000000), Emu(20000000), Emu(800000), [
    "카테고리 간 DOI 중앙값 차이가 크지 않음 → DOI는 카테고리보다 개별 제품 특성에 의해 결정",
], FONT_NAME, Pt(22), COLOR_GRAY)


# ─── SLIDE S4-5: 데이터 누수 진단 ───
print("  Adding S4-5: 데이터 누수 진단...")
s4_5 = create_new_slide()
setup_slide(s4_5)
add_slide_title(s4_5, "데이터 누수(Data Leakage) 진단")

add_multiline_textbox(s4_5, Emu(1905000), Emu(3500000), Emu(9500000), Emu(4500000), [
    "DOI = QOH / ADS  (97.2% 일치 확인)",
    "→ QOH, ADS를 피처로 사용 시 확정 누수",
    "",
    "5가지 시나리오 비교:",
    "S1 (QOH+ADS)       R²=0.91  ← 누수! 나눗셈 역산",
    "S2 (QOH만)           R²=0.59  ← 부분 누수",
    "S3 (ADS만)           R²=-0.08  ← 안전",
    "S4 (둘 다 제거)      R²=-0.07  ← 안전, 정보 부족",
    "S5 (S4+Cat+ABC)  R²=0.25  ← 채택",
], FONT_NAME, Pt(22), COLOR_DARK)

img_path = os.path.join(IMG_DIR, "leakage_scenario_comparison_s4.png")
if os.path.exists(img_path):
    s4_5.shapes.add_picture(img_path, Emu(12000000), Emu(3400000), Emu(11000000), Emu(8000000))

add_textbox(s4_5, Emu(1905000), Emu(11800000), Emu(20000000), Emu(800000),
    "S5 채택: 누수 없이 확보 가능한 최대 예측력(R²=0.25), 잔존 누수 검증 완료 (모든 단일 피처 R²<0.09)",
    FONT_NAME, Pt(20), COLOR_GRAY)


# ─── SLIDE S4-6: Phase A 결과 (모델 비교) ───
print("  Adding S4-6: Phase A 모델 비교...")
s4_6 = create_new_slide()
setup_slide(s4_6)
add_slide_title(s4_6, "Phase A: 회귀 모델 비교 (Default vs Tuned)")

img_path = os.path.join(IMG_DIR, "model_comparison_phaseA_s4.png")
if os.path.exists(img_path):
    s4_6.shapes.add_picture(img_path, Emu(1200000), Emu(3400000), Emu(11000000), Emu(7500000))

img_path2 = os.path.join(IMG_DIR, "default_vs_tuned_phaseA_s4.png")
if os.path.exists(img_path2):
    s4_6.shapes.add_picture(img_path2, Emu(12500000), Emu(3400000), Emu(11000000), Emu(7500000))

add_multiline_textbox(s4_6, Emu(1200000), Emu(11200000), Emu(22000000), Emu(1500000), [
    "Enhanced XGB Tuned: Test R²=0.42, Gap=0.33 (최적)  |  LR: Gap=0.05 (안정적, 낮은 성능)",
    "낮은 R²는 모델 실패가 아닌 누수 대응(QOH+ADS 제거)의 당연한 결과",
], FONT_NAME, Pt(20), COLOR_GRAY)


# ─── SLIDE S4-7: Phase A Feature Importance ───
print("  Adding S4-7: Feature Importance...")
s4_7 = create_new_slide()
setup_slide(s4_7)
add_slide_title(s4_7, "Phase A: Feature Importance 3중 교차검증")

img_path = os.path.join(IMG_DIR, "feature_importance_doi_s4.png")
if os.path.exists(img_path):
    s4_7.shapes.add_picture(img_path, Emu(1200000), Emu(3400000), Emu(14000000), Emu(7500000))

add_multiline_textbox(s4_7, Emu(15500000), Emu(3600000), Emu(8000000), Emu(7000000), [
    "3중 교차검증 결과",
    "",
    "1위: ABC_Class_C",
    "    (3개 방법 모두 1위)",
    "",
    "2위: Reorder_Urgency",
    "    (파생변수 중 최고 중요도)",
    "",
    "3위: Unit_Cost_USD",
    "",
    "→ ABC 등급, 재주문 긴급도,",
    "   단가, 리드타임이",
    "   DOI의 핵심 결정 요인",
], FONT_NAME, Pt(22), COLOR_DARK)


# ─── SLIDE S4-8: Phase A 실제 vs 예측 ───
print("  Adding S4-8: 실제 vs 예측...")
s4_8 = create_new_slide()
setup_slide(s4_8)
add_slide_title(s4_8, "Phase A: 실제 DOI vs 예측 DOI")

img_path = os.path.join(IMG_DIR, "actual_vs_predicted_doi_s4.png")
if os.path.exists(img_path):
    s4_8.shapes.add_picture(img_path, Emu(4500000), Emu(3200000), Emu(15000000), Emu(8500000))

add_multiline_textbox(s4_8, Emu(1905000), Emu(11800000), Emu(20000000), Emu(1000000), [
    "XGB Tuned R²=0.42 | 잔차 정규성 만족 (Shapiro-Wilk p=0.18) | Phase A의 핵심 가치는 피처 해석에 있음",
], FONT_NAME, Pt(20), COLOR_GRAY)


# ─── SLIDE S4-9: Phase B 클러스터링 사전 검증 ───
print("  Adding S4-9: Phase B 사전 검증...")
s4_9 = create_new_slide()
setup_slide(s4_9)
add_slide_title(s4_9, "Phase B: 클러스터링 사전 검증 (Hopkins, Elbow, Silhouette)")

img_path = os.path.join(IMG_DIR, "elbow_method_s4.png")
if os.path.exists(img_path):
    s4_9.shapes.add_picture(img_path, Emu(1200000), Emu(3400000), Emu(11000000), Emu(7500000))

img_path2 = os.path.join(IMG_DIR, "silhouette_scores_s4.png")
if os.path.exists(img_path2):
    s4_9.shapes.add_picture(img_path2, Emu(12500000), Emu(3400000), Emu(11000000), Emu(7500000))

add_multiline_textbox(s4_9, Emu(1200000), Emu(11200000), Emu(22000000), Emu(1500000), [
    "Hopkins Statistic: 0.86 (강한 군집 경향)  |  K=2 최적 (Elbow + Silhouette 모두 지지)",
    "Silhouette Score: Baseline 0.3848 / Enhanced 0.3561  |  안정성 std=0.0009",
], FONT_NAME, Pt(20), COLOR_GRAY)


# ─── SLIDE S4-10: PCA & 군집 특성 ───
print("  Adding S4-10: PCA & 군집 특성...")
s4_10 = create_new_slide()
setup_slide(s4_10)
add_slide_title(s4_10, "Phase B: PCA 2D 시각화 & 군집별 Z-Score 비교")

img_path = os.path.join(IMG_DIR, "pca_clusters_s4.png")
if os.path.exists(img_path):
    s4_10.shapes.add_picture(img_path, Emu(1200000), Emu(3400000), Emu(11000000), Emu(7500000))

img_path2 = os.path.join(IMG_DIR, "cluster_feature_comparison_s4.png")
if os.path.exists(img_path2):
    s4_10.shapes.add_picture(img_path2, Emu(12500000), Emu(3400000), Emu(11000000), Emu(7500000))

add_multiline_textbox(s4_10, Emu(1200000), Emu(11200000), Emu(22000000), Emu(1500000), [
    "Cluster 0 (818개, 82%): 일반 제품군  |  Cluster 1 (182개, 18%): 고관리 필요 제품군",
    "Cluster 1: Reorder_Point z=+1.63, Safety_Stock z=+1.63, Supply_Risk z=+0.88, Available_Stock z=-0.76",
], FONT_NAME, Pt(20), COLOR_GRAY)


# ─── SLIDE S4-11: 레이더 차트 & 군집 프로필 ───
print("  Adding S4-11: 레이더 차트...")
s4_11 = create_new_slide()
setup_slide(s4_11)
add_slide_title(s4_11, "Phase B: 군집별 프로필 비교")

img_path = os.path.join(IMG_DIR, "radar_chart_s4.png")
if os.path.exists(img_path):
    s4_11.shapes.add_picture(img_path, Emu(1200000), Emu(3400000), Emu(11500000), Emu(8000000))

add_multiline_textbox(s4_11, Emu(13500000), Emu(3600000), Emu(9500000), Emu(7500000), [
    "Cluster 0 (일반 제품, 82%)",
    "• Available_Stock: 양수 (+103.6)",
    "• Low Stock 비율: 17.1%",
    "• 평균 EOQ: 974개",
    "• 전략: 정기 발주, 과재고 방지",
    "",
    "Cluster 1 (고관리 제품, 18%)",
    "• Available_Stock: 음수 (-102.3)",
    "• Low Stock 비율: 55.5%",
    "• 음수 Available_Stock: 64.8%",
    "• 평균 EOQ: 2,129개",
    "• 전략: 안전재고 강화 + 긴급 보충",
    "",
    "→ A등급(고매출) 비율이 높아",
    "  비즈니스 영향도가 큰 군집",
], FONT_NAME, Pt(22), COLOR_DARK)


# ─── SLIDE S4-12: ABC & 카테고리 분포 ───
print("  Adding S4-12: ABC & 카테고리 분포...")
s4_12 = create_new_slide()
setup_slide(s4_12)
add_slide_title(s4_12, "Phase B: ABC 등급 & 카테고리별 클러스터 분포")

img_path = os.path.join(IMG_DIR, "abc_cluster_distribution_s4.png")
if os.path.exists(img_path):
    s4_12.shapes.add_picture(img_path, Emu(1200000), Emu(3400000), Emu(11000000), Emu(7500000))

img_path2 = os.path.join(IMG_DIR, "category_cluster_distribution_s4.png")
if os.path.exists(img_path2):
    s4_12.shapes.add_picture(img_path2, Emu(12500000), Emu(3400000), Emu(11000000), Emu(7500000))

add_multiline_textbox(s4_12, Emu(1200000), Emu(11200000), Emu(22000000), Emu(1500000), [
    "Cluster 1에 A등급(고매출) 비율 높음 → 비즈니스 영향 큰 핵심 제품이 고위험 군집에 속함",
    "Beverages(음료), Fresh Produce(신선) 카테고리에서 Cluster 1 비율 상대적으로 높음",
], FONT_NAME, Pt(20), COLOR_GRAY)


# ─── SLIDE S4-13: Phase C EOQ 시뮬레이션 ───
print("  Adding S4-13: Phase C EOQ...")
s4_13 = create_new_slide()
setup_slide(s4_13)
add_slide_title(s4_13, "Phase C: 군집별 EOQ 발주 전략")

add_multiline_textbox(s4_13, Emu(1905000), Emu(3500000), Emu(9500000), Emu(8500000), [
    "EOQ = sqrt(2 x D x S / H)",
    "  D: 연간 수요량",
    "  S: 1회 발주 비용 = 50 USD",
    "  H: 단위당 연간 보관비 = Unit_Cost x 20%",
    "",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "",
    "Cluster 0 (일반 제품)",
    "• 전략: EOQ 기반 정기 발주",
    "• 평균 EOQ: 974개",
    "• 핵심: 과재고 방지",
    "",
    "Cluster 1 (고관리 제품)",
    "• 전략: 안전재고 강화 + 즉시 보충",
    "• 평균 EOQ: 2,129개",
    "• 핵심: 결품 방지 (Available_Stock 음수 해소)",
    "• 추가: 안전재고 상향, 대체 공급처 확보",
], FONT_NAME, Pt(24), COLOR_DARK)

add_multiline_textbox(s4_13, Emu(12500000), Emu(3500000), Emu(10500000), Emu(8500000), [
    "EOQ 민감도 분석",
    "",
    "발주비용(S)    C0 EOQ    C1 EOQ",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "$25                ~689       ~1,505",
    "$50 (기본)       974        2,129",
    "$75                ~1,193    ~2,607",
    "$100              ~1,378    ~3,011",
    "",
    "S가 2배 → EOQ ~1.41배 증가",
    "(EOQ 공식에서 S는 제곱근 안에 위치)",
    "",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "Cluster 1의 EOQ가 2배 이상 높음",
    "→ 연간 수요가 크고 한 번에 더 많이",
    "  발주해야 경제적",
    "→ 64.8% 음수 재고 → EOQ 도달 전",
    "  긴급 발주 필요",
], FONT_NAME, Pt(22), COLOR_DARK)


# ─── SLIDE S4-14: 소주제 연결 인사이트 ───
print("  Adding S4-14: 소주제 연결 인사이트...")
s4_14 = create_new_slide()
setup_slide(s4_14)
add_slide_title(s4_14, "소주제 1~4 종합 연결 인사이트")

add_multiline_textbox(s4_14, Emu(1905000), Emu(3500000), Emu(20000000), Emu(8500000), [
    "소주제 1 연계  |  Cluster 1: Low Stock 55.5% → 재고 부족 경보 시스템 필요",
    "                           재고 상태 분류 결과를 군집별 모니터링에 활용",
    "",
    "소주제 2 연계  |  Cluster 1: 고 ADS(높은 일일 판매량)",
    "                           판매량 예측 모델 결과를 발주 타이밍 결정에 활용",
    "",
    "소주제 3 연계  |  Cluster 0: Available_Stock 양수 → 과재고 위험",
    "                           폐기 위험 예측과 연계하여 과재고 방지",
    "",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "",
    "4개 소주제의 통합 의사결정 사이클:",
    "",
    "  소주제 1 (상태 진단)  →  소주제 2 (수요 예측)  →  소주제 3 (위험 감지)  →  소주제 4 (전략 수립)",
    "  현재 상태 파악              미래 수요 예측              위험 사전 식별              최적 행동 결정",
    "",
    "→ 식료품 유통 재고 관리의 완전한 의사결정 체계 완성",
], FONT_NAME, Pt(24), COLOR_DARK)


# ─── SLIDE S4-15: 결론 & 한계점 ───
print("  Adding S4-15: 결론 & 한계점...")
s4_15 = create_new_slide()
setup_slide(s4_15)
add_slide_title(s4_15, "소주제 4: 결론 및 한계점")

add_multiline_textbox(s4_15, Emu(1905000), Emu(3500000), Emu(10000000), Emu(8500000), [
    "핵심 발견 5가지",
    "",
    "1. Phase A: 누수 대응 후에도 ABC 등급,",
    "   재주문 긴급도, 단가, 리드타임이",
    "   DOI의 핵심 결정 요인임을 3중 검증",
    "",
    "2. Phase B: Hopkins(0.86) 확인 후",
    "   K=2 최적, 안정성 std=0.0009",
    "",
    "3. Cluster 1(18%)은 '고위험/고관리'",
    "   제품군. A등급 비율 높아 비즈니스",
    "   영향도가 큼",
    "",
    "4. EOQ로 구체적 발주량 제시:",
    "   C0=974(정기), C1=2,129(긴급)",
    "",
    "5. 4개 소주제 통합 재고 관리",
    "   의사결정 사이클 완성",
], FONT_NAME, Pt(22), COLOR_DARK)

add_multiline_textbox(s4_15, Emu(12500000), Emu(3500000), Emu(10500000), Emu(8500000), [
    "한계점",
    "",
    "데이터 한계",
    "• 1,000건 규모 (대규모 유통 대비 표본 작음)",
    "• 시계열 정보 부재 (스냅샷 데이터)",
    "• 수요 변동성 부재",
    "",
    "Phase A 한계",
    "• 낮은 R²(0.42) — 구조적 한계",
    "• 트리 모델 과적합 유지 (Gap 0.33~0.47)",
    "",
    "Phase B 한계",
    "• K-Means 구형 가정",
    "• K=2의 단순성",
    "• DBSCAN Noise ~15%",
    "",
    "Phase C 한계",
    "• EOQ 가정의 단순성 (일정 수요,",
    "  고정 비용, 즉시 보충)",
], FONT_NAME, Pt(22), COLOR_DARK)


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════

print(f"\nTotal slides after additions: {len(prs.slides)}")
print(f"Saving to: {PPTX_PATH}")
prs.save(PPTX_PATH)
print("DONE!")
